#!/usr/bin/env python3
"""Fetch a Google Drive document and print its text content.

Usage:
    python fetch_doc.py <google_drive_url_or_doc_id>

Requires env vars: GOOGLE_CLIENT_ID, GOOGLE_CLIENT_SECRET, GOOGLE_REFRESH_TOKEN
"""

import os
import re
import sys
from urllib.parse import parse_qs, urlparse


def extract_doc_id(url_or_id: str) -> str:
    """Extract document ID from a Google Drive URL or return as-is if already an ID."""
    # If it looks like a raw ID (no slashes, no dots)
    if re.match(r'^[a-zA-Z0-9_-]+$', url_or_id):
        return url_or_id
    # Try /d/{ID}/ pattern
    m = re.search(r'/d/([a-zA-Z0-9_-]+)', url_or_id)
    if m:
        return m.group(1)
    # Try ?id={ID} pattern
    parsed = urlparse(url_or_id)
    qs = parse_qs(parsed.query)
    if 'id' in qs:
        return qs['id'][0]
    raise ValueError(f"Cannot extract document ID from: {url_or_id}")


def main():
    # Ensure UTF-8 output on Windows (cp1252 can't handle Unicode chars in docs)
    if hasattr(sys.stdout, 'reconfigure'):
        sys.stdout.reconfigure(encoding='utf-8', errors='replace')
    if hasattr(sys.stderr, 'reconfigure'):
        sys.stderr.reconfigure(encoding='utf-8', errors='replace')

    if len(sys.argv) < 2:
        print("Usage: python fetch_doc.py <google_drive_url_or_doc_id>", file=sys.stderr)
        sys.exit(1)

    url = sys.argv[1]

    # Check env vars
    for var in ('GOOGLE_CLIENT_ID', 'GOOGLE_CLIENT_SECRET', 'GOOGLE_REFRESH_TOKEN'):
        if not os.environ.get(var):
            print(f"ERROR: Missing environment variable {var}", file=sys.stderr)
            sys.exit(1)

    try:
        from google.oauth2.credentials import Credentials
        from googleapiclient.discovery import build
    except ImportError:
        # Auto-install
        import subprocess
        subprocess.check_call([sys.executable, '-m', 'pip', 'install', '-q',
                               'google-api-python-client', 'google-auth'])
        from google.oauth2.credentials import Credentials
        from googleapiclient.discovery import build

    doc_id = extract_doc_id(url)

    creds = Credentials(
        token=None,
        refresh_token=os.environ['GOOGLE_REFRESH_TOKEN'],
        token_uri='https://oauth2.googleapis.com/token',
        client_id=os.environ['GOOGLE_CLIENT_ID'],
        client_secret=os.environ['GOOGLE_CLIENT_SECRET'],
    )
    service = build('drive', 'v3', credentials=creds)

    # Get metadata
    try:
        meta = service.files().get(fileId=doc_id, fields='mimeType,name').execute()
    except Exception as e:
        err = str(e)
        if '404' in err:
            print(f"ERROR: Cannot access document (ID: {doc_id}).")
            print("The OAuth account may not have access to this file.")
            print("Possible fixes: share the doc with the authorized Google account,")
            print("or re-run scripts/setup_google_auth.py with the correct account.")
        else:
            print(f"ERROR: Google API error: {err[:300]}")
        sys.exit(1)

    mime = meta['mimeType']
    name = meta.get('name', 'untitled')

    print(f"=== {name} ===")
    print(f"Type: {mime}\n")

    text = None

    if mime == 'application/vnd.google-apps.document':
        content = service.files().export(fileId=doc_id, mimeType='text/plain').execute()
        text = content.decode('utf-8') if isinstance(content, bytes) else content
    elif mime == 'application/vnd.google-apps.spreadsheet':
        content = service.files().export(fileId=doc_id, mimeType='text/csv').execute()
        text = content.decode('utf-8') if isinstance(content, bytes) else content
    elif mime == 'application/vnd.google-apps.presentation':
        content = service.files().export(fileId=doc_id, mimeType='text/plain').execute()
        text = content.decode('utf-8') if isinstance(content, bytes) else content
    elif mime in (
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'application/msword',
    ):
        # Download .docx/.doc and extract text
        import io
        import tempfile
        raw = service.files().get_media(fileId=doc_id).execute()
        try:
            import docx
        except ImportError:
            import subprocess as _sp
            _sp.check_call([sys.executable, '-m', 'pip', 'install', '-q', 'python-docx'])
            import docx
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(raw)
            tmp_path = tmp.name
        try:
            doc = docx.Document(tmp_path)
            text = '\n'.join(p.text for p in doc.paragraphs)
        finally:
            os.unlink(tmp_path)
    elif mime in (
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        'application/vnd.ms-excel',
    ):
        # Download .xlsx/.xls and convert to CSV text
        raw = service.files().get_media(fileId=doc_id).execute()
        try:
            import openpyxl
        except ImportError:
            import subprocess as _sp
            _sp.check_call([sys.executable, '-m', 'pip', 'install', '-q', 'openpyxl'])
            import openpyxl
        import io
        wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"--- Sheet: {sheet_name} ---")
            for row in ws.iter_rows(values_only=True):
                parts.append(','.join(str(c) if c is not None else '' for c in row))
        wb.close()
        text = '\n'.join(parts)
    elif mime in (
        'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        'application/vnd.ms-powerpoint',
    ):
        # Download .pptx/.ppt and extract text
        raw = service.files().get_media(fileId=doc_id).execute()
        try:
            from pptx import Presentation
        except ImportError:
            import subprocess as _sp
            _sp.check_call([sys.executable, '-m', 'pip', 'install', '-q', 'python-pptx'])
            from pptx import Presentation
        import io
        prs = Presentation(io.BytesIO(raw))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"--- Slide {i} ---")
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    parts.append(shape.text)
        text = '\n'.join(parts)
    elif mime == 'application/pdf':
        # Download PDF and save to workspace for further processing
        raw = service.files().get_media(fileId=doc_id).execute()
        workspace = os.environ.get('WORKSPACE_PATH', '/mnt/user-data/workspace')
        safe_name = re.sub(r'[^\w\-.]', '_', name)
        save_path = os.path.join(workspace, safe_name)
        with open(save_path, 'wb') as f:
            f.write(raw)
        print(f"Downloaded PDF to: {save_path}")
        print(f"Size: {len(raw)} bytes")
        print("Use Python to process this file further.")
        sys.exit(0)
    else:
        # Unknown binary type — download to workspace so the agent can process it
        raw = service.files().get_media(fileId=doc_id).execute()
        workspace = os.environ.get('WORKSPACE_PATH', '/mnt/user-data/workspace')
        safe_name = re.sub(r'[^\w\-.]', '_', name)
        save_path = os.path.join(workspace, safe_name)
        with open(save_path, 'wb') as f:
            f.write(raw)
        print(f"Binary file type ({mime}). Downloaded to: {save_path}")
        print(f"Size: {len(raw)} bytes")
        print("Use Python to process this file further.")
        sys.exit(0)

    if text:
        print(text)


if __name__ == '__main__':
    main()
